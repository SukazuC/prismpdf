import logging
from pathlib import Path


logger = logging.getLogger(__name__)


def convert_to_pptx(input_path: Path, output_path: Path):
    """
    Convert PDF to PPTX using PyMuPDF for page rendering + python-pptx for slides.
    Each page becomes a slide with the rendered image as background.
    Result is a visual slide deck (non-editable text).
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        import io

        doc = fitz.open(str(input_path))
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)

            # Calculate scale to fit slide dimensions
            page_width = page.rect.width
            page_height = page.rect.height
            scale_x = 13.333 / (page_width / 72)
            scale_y = 7.5 / (page_height / 72)
            scale = min(scale_x, scale_y) * 0.9  # 90% to leave margins

            render_scale = scale * 2  # 2x for retina quality
            mat = fitz.Matrix(render_scale, render_scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            img_data = pix.tobytes("png")

            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add image to slide
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(
                img_stream,
                Inches(0.5),
                Inches(0.5),
                width=Inches(12.333),
                height=Inches(6.5),
            )

        page_count = len(doc)
        doc.close()
        prs.save(str(output_path))

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("PPTX output is empty")

        logger.info(
            "PPTX conversion successful: %s -> %s (%d pages, %d bytes)",
            input_path.name,
            output_path.name,
            page_count,
            output_path.stat().st_size,
        )

    except Exception as e:
        logger.error("PPTX conversion failed: %s", str(e))
        raise
